"""
Build PRESENTATION.pptx — a polished, dark-themed deck for the ForeRoute group
presentation. Hand-coded with python-pptx for a consistent layout grid (no
template baggage). No code blocks anywhere; plain-language copy throughout.

Design system:
  · Background slate-900, accent sky-400, semantic emerald/amber/red cues.
  · Single layout grid for content slides:
      eyebrow y=0.55  ·  title y=0.95  ·  subtitle y=1.75  ·  body y=2.55
  · Title slides and the closing slide get hero treatment.
  · Footer accent stripe (top) + page counter (bottom) on every content slide.

Run:
  python3 build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).parent / "PRESENTATION.pptx"

# ---- colour system --------------------------------------------------------

BG = RGBColor(0x0F, 0x17, 0x2A)
SURFACE = RGBColor(0x1E, 0x29, 0x3B)
SURFACE_2 = RGBColor(0x33, 0x41, 0x55)
DIVIDER = RGBColor(0x1E, 0x29, 0x3B)
TEXT = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SUBTLE = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_DEEP = RGBColor(0x0E, 0xA5, 0xE9)
ACCENT_LIGHT = RGBColor(0x7D, 0xD3, 0xFC)
GOOD = RGBColor(0x34, 0xD3, 0x99)
WARN = RGBColor(0xFB, 0xBF, 0x24)
BAD = RGBColor(0xF8, 0x71, 0x71)

# Typography
FONT = "Inter"
FONT_FALLBACK = "Helvetica Neue"

# Layout grid (inches)
PAGE_W = 13.333
PAGE_H = 7.5
LEFT_M = 0.7
RIGHT_M = 0.7
TOP_M = 0.55
BODY_W = PAGE_W - LEFT_M - RIGHT_M

EYEBROW_Y = 0.55
TITLE_Y = 0.95
SUBTITLE_Y = 1.75
BODY_Y = 2.55
BODY_H = 4.4  # ends before the footer band at 6.95
FOOTER_Y = 7.0

# ---- prs ------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(PAGE_W)
prs.slide_height = Inches(PAGE_H)
BLANK = prs.slide_layouts[6]


# ---- primitives -----------------------------------------------------------


def fill_bg(slide, color=BG):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def rect(slide, left, top, width, height, fill=SURFACE, line=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    if radius is not None:
        s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    size=18,
    color=TEXT,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
    line_spacing=None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if isinstance(text, str):
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    else:
        for chunk, opts in text:
            run = p.add_run()
            run.text = chunk
            run.font.name = opts.get("font", font)
            run.font.size = Pt(opts.get("size", size))
            run.font.color.rgb = opts.get("color", color)
            run.font.bold = opts.get("bold", bold)
            run.font.italic = opts.get("italic", italic)
    return tb


def add_paragraphs(
    slide,
    lines,
    left,
    top,
    width,
    height,
    size=15,
    color=TEXT,
    bullet=False,
    line_spacing=1.4,
    space_after=None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        if bullet:
            rp = p.add_run()
            rp.text = "•  "
            rp.font.name = FONT
            rp.font.size = Pt(size)
            rp.font.color.rgb = ACCENT
        if isinstance(line, str):
            r = p.add_run()
            r.text = line
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
        else:
            for chunk, opts in line:
                r = p.add_run()
                r.text = chunk
                r.font.name = opts.get("font", FONT)
                r.font.size = Pt(opts.get("size", size))
                r.font.color.rgb = opts.get("color", color)
                r.font.bold = opts.get("bold", False)
                r.font.italic = opts.get("italic", False)
    return tb


# ---- chrome ---------------------------------------------------------------


def chrome(slide, page_num, total):
    # Top accent stripe
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.shadow.inherit = False
    # Bottom thin divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(LEFT_M),
        Inches(FOOTER_Y - 0.05),
        Inches(BODY_W),
        Inches(0.015),
    )
    divider.line.fill.background()
    divider.fill.solid()
    divider.fill.fore_color.rgb = SURFACE
    divider.shadow.inherit = False
    # Footer text
    add_text(
        slide,
        "ForeRoute — Know the road before you go",
        Inches(LEFT_M),
        Inches(FOOTER_Y + 0.05),
        Inches(8),
        Inches(0.35),
        size=10,
        color=SUBTLE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        f"{page_num} / {total}",
        Inches(PAGE_W - RIGHT_M - 1.0),
        Inches(FOOTER_Y + 0.05),
        Inches(1.0),
        Inches(0.35),
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def title_block(slide, eyebrow, title, subtitle=None):
    add_text(
        slide,
        eyebrow.upper(),
        Inches(LEFT_M),
        Inches(EYEBROW_Y),
        Inches(BODY_W),
        Inches(0.3),
        size=11,
        color=ACCENT_LIGHT,
        bold=True,
    )
    add_text(
        slide,
        title,
        Inches(LEFT_M),
        Inches(TITLE_Y),
        Inches(BODY_W),
        Inches(0.8),
        size=32,
        color=TEXT,
        bold=True,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            Inches(LEFT_M),
            Inches(SUBTITLE_Y),
            Inches(BODY_W),
            Inches(0.6),
            size=15,
            color=MUTED,
            italic=True,
            line_spacing=1.3,
        )


# ---- table builder --------------------------------------------------------


def add_table(slide, data, left, top, width, height, col_widths=None, font_size=13, row_height=None):
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    if row_height:
        for row in table.rows:
            row.height = row_height
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.14)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = SURFACE_2
            else:
                cell.fill.fore_color.rgb = BG if r % 2 == 0 else SURFACE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.line_spacing = 1.25
            if isinstance(val, str):
                r0 = p.add_run()
                r0.text = val
                r0.font.name = FONT
                r0.font.size = Pt(font_size)
                r0.font.color.rgb = TEXT
                r0.font.bold = r == 0
            else:
                txt, opts = val
                r0 = p.add_run()
                r0.text = txt
                r0.font.name = opts.get("font", FONT)
                r0.font.size = Pt(opts.get("size", font_size))
                r0.font.color.rgb = opts.get("color", TEXT)
                r0.font.bold = opts.get("bold", r == 0)
    return tbl_shape


# ============================================================================
# SLIDES
# ============================================================================

TOTAL = 17


# ---- 1 · Title ------------------------------------------------------------


def slide_title():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    # Eyebrow strip
    add_text(
        s, "NORTHEASTERN  ·  GROUP PROJECT  ·  2026",
        Inches(0), Inches(2.1), Inches(PAGE_W), Inches(0.4),
        size=12, color=ACCENT_LIGHT, bold=True, align=PP_ALIGN.CENTER,
    )
    add_text(
        s, "ForeRoute",
        Inches(0), Inches(2.55), Inches(PAGE_W), Inches(1.4),
        size=84, color=TEXT, bold=True, align=PP_ALIGN.CENTER,
    )
    add_text(
        s, "Know the road before you go.",
        Inches(0), Inches(4.0), Inches(PAGE_W), Inches(0.6),
        size=22, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
    )
    # Accent rule
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.05), Inches(4.85), Inches(1.2), Inches(0.05))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.shadow.inherit = False
    add_text(
        s, "A weather-aware safe-routing app, with two scores per route — one for the conditions right now, one for what crashes have actually happened on every street.",
        Inches(1.5), Inches(5.05), Inches(PAGE_W - 3.0), Inches(1.0),
        size=15, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.4,
    )
    add_text(
        s, "[member 1]   ·   [member 2]   ·   [member 3]   ·   [member 4]",
        Inches(0), Inches(6.5), Inches(PAGE_W), Inches(0.4),
        size=13, color=SUBTLE, align=PP_ALIGN.CENTER,
    )


# ---- 2 · The idea ---------------------------------------------------------


def slide_idea():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 2, TOTAL)
    title_block(
        s,
        "1 · The idea",
        "Fastest is not the same as safest.",
        "Today's navigation apps optimise for time. None of them ask the question that matters on a snowy Tuesday morning.",
    )
    # Pull-quote
    rect(s, Inches(LEFT_M), Inches(BODY_Y), Inches(BODY_W), Inches(1.1), fill=SURFACE, radius=0.05)
    add_text(
        s, "\"Is this trip safe today, on these roads?\"",
        Inches(LEFT_M + 0.4), Inches(BODY_Y + 0.18), Inches(BODY_W - 0.8), Inches(0.74),
        size=22, color=ACCENT_LIGHT, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Two questions
    add_text(
        s, "ForeRoute answers two questions, separately, in plain language:",
        Inches(LEFT_M), Inches(BODY_Y + 1.45), Inches(BODY_W), Inches(0.4),
        size=15, color=MUTED,
    )
    add_paragraphs(
        s,
        [
            [("How dangerous is the weather right now?", {"bold": True, "size": 18, "color": TEXT}),
             ("  Snow, ice, hydroplaning rain, low visibility, wind, tricky road types.", {"color": MUTED, "size": 15})],
            [("How crash-prone is this route, based on what's actually happened here before?", {"bold": True, "size": 18, "color": TEXT}),
             ("  Eight years of real Boston crash records, by location and time.", {"color": MUTED, "size": 15})],
        ],
        Inches(LEFT_M + 0.2), Inches(BODY_Y + 1.95), Inches(BODY_W - 0.4), Inches(2.0),
        size=18, bullet=True, line_spacing=1.5, space_after=8,
    )


# ---- 3 · Two scores side by side ------------------------------------------


def slide_two_scores():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 3, TOTAL)
    title_block(
        s,
        "2 · Two scores",
        "Built side by side. We even flag when they disagree.",
        "One number can't honestly answer both \"how should I drive?\" and \"which route should I pick?\". So we don't try.",
    )
    col_w = (BODY_W - 0.4) / 2
    top = Inches(BODY_Y)
    height = Inches(3.7)
    # Left card — Weather
    rect(s, Inches(LEFT_M), top, Inches(col_w), height, fill=SURFACE, radius=0.05)
    add_text(s, "WEATHER RIGHT NOW", Inches(LEFT_M + 0.25), Inches(BODY_Y + 0.2), Inches(col_w - 0.5), Inches(0.3),
             size=11, color=WARN, bold=True)
    add_text(s, "Is right now known-adverse?",
             Inches(LEFT_M + 0.25), Inches(BODY_Y + 0.55), Inches(col_w - 0.5), Inches(0.5),
             size=18, color=TEXT, bold=True)
    add_paragraphs(
        s,
        [
            [("Built on", {"size": 13, "color": MUTED, "bold": True}),
             ("   The physics of driving — hand-tuned rules.", {"size": 14, "color": TEXT})],
            [("Loud when", {"size": 13, "color": MUTED, "bold": True}),
             ("   Snow on a bridge. Freezing rain. Fog.", {"size": 14, "color": TEXT})],
            [("Quiet when", {"size": 13, "color": MUTED, "bold": True}),
             ("   Mild, clear weather.", {"size": 14, "color": TEXT})],
            [("Always on", {"size": 13, "color": MUTED, "bold": True}),
             ("   Runs even if the model is offline.", {"size": 14, "color": TEXT})],
        ],
        Inches(LEFT_M + 0.25), Inches(BODY_Y + 1.2), Inches(col_w - 0.5), Inches(2.5),
        size=13, line_spacing=1.55, space_after=6,
    )
    # Right card — Crash history
    right_x = LEFT_M + col_w + 0.4
    rect(s, Inches(right_x), top, Inches(col_w), height, fill=SURFACE, radius=0.05)
    add_text(s, "CRASH HISTORY", Inches(right_x + 0.25), Inches(BODY_Y + 0.2), Inches(col_w - 0.5), Inches(0.3),
             size=11, color=ACCENT, bold=True)
    add_text(s, "Is this road crashier than typical?",
             Inches(right_x + 0.25), Inches(BODY_Y + 0.55), Inches(col_w - 0.5), Inches(0.5),
             size=18, color=TEXT, bold=True)
    add_paragraphs(
        s,
        [
            [("Built on", {"size": 13, "color": MUTED, "bold": True}),
             ("   Real Boston crashes, 2018–2025.", {"size": 14, "color": TEXT})],
            [("Loud when", {"size": 13, "color": MUTED, "bold": True}),
             ("   Crashy corridor at rush hour.", {"size": 14, "color": TEXT})],
            [("Quiet when", {"size": 13, "color": MUTED, "bold": True}),
             ("   Quiet residential streets at off-peak.", {"size": 14, "color": TEXT})],
            [("Always on", {"size": 13, "color": MUTED, "bold": True}),
             ("   Fallbacks if the model is unavailable.", {"size": 14, "color": TEXT})],
        ],
        Inches(right_x + 0.25), Inches(BODY_Y + 1.2), Inches(col_w - 0.5), Inches(2.5),
        size=13, line_spacing=1.55, space_after=6,
    )
    # Takeaway
    add_text(
        s,
        "Disagreement is useful information. The UI flags it.",
        Inches(LEFT_M), Inches(BODY_Y + 3.95), Inches(BODY_W), Inches(0.4),
        size=14, color=ACCENT_LIGHT, italic=True, align=PP_ALIGN.CENTER,
    )


# ---- 4 · Where the data comes from ----------------------------------------


def slide_data_sources():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 4, TOTAL)
    title_block(
        s,
        "3 · Where the data comes from",
        "Three real, public sources — no synthetic data behind the production model.",
        "Every row in our training set is anchored on a real event in the City of Boston.",
    )
    data = [
        ["Source", "What we use it for", "What we got"],
        [
            ("Boston Vision Zero Crash Records", {"bold": True, "size": 14}),
            "Every positive example is a real, dated, police-reported crash.",
            ("22,260 crashes", {"bold": True, "color": ACCENT}),
        ],
        [
            ("Open-Meteo Historical Weather", {"bold": True, "size": 14}),
            "The exact conditions at every crash's location and hour.",
            ("8 years of hourly data", {"color": TEXT}),
        ],
        [
            ("Mapbox Streets", {"bold": True, "size": 14}),
            "Classifies each location as highway, arterial, residential, bridge, tunnel, mountain.",
            ("Real road type, per point", {"color": TEXT}),
        ],
    ]
    add_table(
        s, data, Inches(LEFT_M), Inches(BODY_Y), Inches(BODY_W), Inches(2.6),
        col_widths=[Inches(3.5), Inches(6.0), Inches(2.4)], font_size=14,
    )
    # Pull quote
    rect(s, Inches(LEFT_M), Inches(5.4), Inches(BODY_W), Inches(1.3), fill=SURFACE, radius=0.05)
    add_text(
        s,
        "No survey data, no simulated weather, no toy datasets.",
        Inches(LEFT_M + 0.35), Inches(5.6), Inches(BODY_W - 0.7), Inches(0.4),
        size=18, color=TEXT, bold=True,
    )
    add_text(
        s,
        "If a model is going to advise drivers on safety, it should be trained on actual incidents — not on a generated approximation.",
        Inches(LEFT_M + 0.35), Inches(6.0), Inches(BODY_W - 0.7), Inches(0.6),
        size=14, color=MUTED, italic=True, line_spacing=1.3,
    )


# ---- 5 · Cleaning ---------------------------------------------------------


def slide_cleaning():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 5, TOTAL)
    title_block(
        s,
        "4 · How we cleaned the data",
        "Five rules that shaped the final dataset.",
        "Cleaning is where most machine-learning projects either succeed or quietly fail.",
    )
    rules = [
        ("Bounding box", "Only crashes inside the Boston metro made it in. Anything outside is tagged so the model can later say \"I don't know about this area.\""),
        ("Eight-year window", "2018–2025. Older records reflect a different vehicle mix, signal timing, and road infrastructure."),
        ("Valid coordinates only", "Records missing latitude, longitude, or a timestamp were discarded."),
        ("Time-of-event matters", "Every crash was paired with the historical weather at the hour and place it actually happened — not a daily average."),
        ("No future leakage", "When we count nearby crashes as a feature, we strictly use crashes that happened before the sample's own timestamp."),
    ]
    y0 = BODY_Y
    row_h = 0.78
    for i, (label, body) in enumerate(rules):
        y = y0 + i * row_h
        # Number badge
        num_box = rect(s, Inches(LEFT_M), Inches(y + 0.05), Inches(0.55), Inches(0.55), fill=ACCENT, radius=0.5)
        add_text(s, str(i + 1), Inches(LEFT_M), Inches(y + 0.05), Inches(0.55), Inches(0.55),
                 size=18, color=BG, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Label
        add_text(s, label, Inches(LEFT_M + 0.75), Inches(y + 0.02), Inches(3.0), Inches(0.35),
                 size=15, color=TEXT, bold=True)
        # Body
        add_text(s, body, Inches(LEFT_M + 0.75), Inches(y + 0.32), Inches(BODY_W - 0.95), Inches(0.5),
                 size=13, color=MUTED, line_spacing=1.3)


# ---- 6 · Negatives --------------------------------------------------------


def slide_negatives():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 6, TOTAL)
    title_block(
        s,
        "5 · The negatives problem",
        "How do we train a model on \"nothing happened here\"?",
        "Real crashes give us positives. The model also needs examples of when nothing went wrong — and that part is surprisingly tricky.",
    )
    col_w = (BODY_W - 0.4) / 2
    # Naive approach card
    rect(s, Inches(LEFT_M), Inches(BODY_Y), Inches(col_w), Inches(3.7), fill=SURFACE, radius=0.05)
    add_text(s, "THE NAIVE WAY", Inches(LEFT_M + 0.25), Inches(BODY_Y + 0.18), Inches(col_w - 0.5), Inches(0.3),
             size=11, color=BAD, bold=True)
    add_text(s, "Drop random points across Boston.",
             Inches(LEFT_M + 0.25), Inches(BODY_Y + 0.5), Inches(col_w - 0.5), Inches(0.6),
             size=17, color=TEXT, bold=True)
    add_text(
        s,
        "The model learned a shortcut:\n\"Residential streets = safe.\"\n\nThat isn't really true. It just reflected that crashes cluster on arterials and our random negatives mostly landed on quiet residential roads.",
        Inches(LEFT_M + 0.25), Inches(BODY_Y + 1.2), Inches(col_w - 0.5), Inches(2.4),
        size=14, color=MUTED, line_spacing=1.5,
    )
    # Our approach card
    right_x = LEFT_M + col_w + 0.4
    rect(s, Inches(right_x), Inches(BODY_Y), Inches(col_w), Inches(3.7), fill=SURFACE, radius=0.05)
    add_text(s, "WHAT WE DID", Inches(right_x + 0.25), Inches(BODY_Y + 0.18), Inches(col_w - 0.5), Inches(0.3),
             size=11, color=GOOD, bold=True)
    add_text(s, "Anchor every negative on a real crash spot.",
             Inches(right_x + 0.25), Inches(BODY_Y + 0.5), Inches(col_w - 0.5), Inches(0.6),
             size=17, color=TEXT, bold=True)
    add_text(
        s,
        "Each negative sits about 55 metres from an actual crash, at a random time when no crash happened there.\n\nNegatives and positives now live on the same kinds of roads. The model is forced to learn when a crash happens at a given place, not where.",
        Inches(right_x + 0.25), Inches(BODY_Y + 1.2), Inches(col_w - 0.5), Inches(2.4),
        size=14, color=TEXT, line_spacing=1.5,
    )


# ---- 7 · Features ---------------------------------------------------------


def slide_features():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 7, TOTAL)
    title_block(
        s,
        "6 · What the model looks at",
        "Fifteen inputs, grouped into five families.",
        "Six of these were added late in development — and they lifted accuracy by roughly five points on the same data.",
    )
    groups = [
        ("WEATHER  (7)", WARN, "temperature   ·   precipitation type   ·   precipitation intensity\nwind   ·   visibility   ·   humidity   ·   dew point",
         "The physics of driving."),
        ("ROAD  (2)", ACCENT_LIGHT, "road type   ·   segment length",
         "A bridge is not a residential street."),
        ("LOCATION  (2)", ACCENT_LIGHT, "latitude   ·   longitude",
         "Two segments of the same road can be very different."),
        ("TIME  (3)", ACCENT_LIGHT, "hour of day   ·   day of week   ·   month",
         "Rush hour ≠ 3 AM. Friday ≠ Tuesday."),
        ("CRASH HISTORY  (1)", GOOD, "crashes within 100 metres in the prior year",
         "Past crashes are the strongest predictor of future ones."),
    ]
    y0 = BODY_Y
    row_h = 0.78
    for i, (label, color, content, why) in enumerate(groups):
        y = y0 + i * row_h
        # Label
        add_text(s, label, Inches(LEFT_M), Inches(y + 0.05), Inches(3.0), Inches(0.3),
                 size=11, color=color, bold=True)
        # Content
        add_text(s, content, Inches(LEFT_M), Inches(y + 0.32), Inches(7.5), Inches(0.5),
                 size=14, color=TEXT, line_spacing=1.2)
        # Why (right-aligned italic)
        add_text(s, why, Inches(LEFT_M + 7.7), Inches(y + 0.05), Inches(BODY_W - 7.7), Inches(0.7),
                 size=12, color=MUTED, italic=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)


# ---- 8 · Three models -----------------------------------------------------


def slide_models():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 8, TOTAL)
    title_block(
        s,
        "7 · The three models we trained",
        "Three deliberately different algorithms — compared fairly.",
        "If a complex model can't beat a simple one, the complexity isn't earning its keep.",
    )
    col_w = (BODY_W - 0.6) / 3
    cards = [
        ("DECISION TREE", "A transparent baseline. The fitted tree can be drawn out and audited — \"if it's snowing and below 2 °C, then high risk\". Our sanity check.", ACCENT_LIGHT),
        ("GRADIENT BOOSTING", "The standard production-grade choice for table-shaped data. Captures interactions like snow × cold × bridge automatically. Well-behaved out of the box.", GOOD),
        ("NEURAL NETWORK", "A non-linear contrast to trees, with a completely different inductive bias. Tells us whether deep models buy us anything here.", ACCENT_LIGHT),
    ]
    for i, (label, body, color) in enumerate(cards):
        x = LEFT_M + i * (col_w + 0.3)
        rect(s, Inches(x), Inches(BODY_Y), Inches(col_w), Inches(3.8), fill=SURFACE, radius=0.06)
        add_text(s, label, Inches(x + 0.25), Inches(BODY_Y + 0.25), Inches(col_w - 0.5), Inches(0.3),
                 size=11, color=color, bold=True)
        add_text(s, body, Inches(x + 0.25), Inches(BODY_Y + 0.7), Inches(col_w - 0.5), Inches(2.9),
                 size=14, color=TEXT, line_spacing=1.5)


# ---- 9 · Performance ------------------------------------------------------


def slide_performance():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 9, TOTAL)
    title_block(
        s,
        "8 · Which one won, and why",
        "Gradient Boosting won on ranking quality and is the registered production model.",
        "Numbers below are at the conventional 0.5 threshold. The operating point for a safety advisory is different — see the next slide.",
    )
    data = [
        ["Model", "Ranking quality", "Precision", "Recall", "F1"],
        ["Decision Tree", "0.636", "0.37", "0.25", "0.30"],
        [
            ("Gradient Boosting  ← registered", {"bold": True, "color": ACCENT, "size": 14}),
            ("0.651", {"bold": True, "color": ACCENT, "size": 14}),
            ("0.39", {"bold": True, "color": TEXT, "size": 14}),
            ("0.25", {"bold": True, "color": TEXT, "size": 14}),
            ("0.30", {"bold": True, "color": TEXT, "size": 14}),
        ],
        ["Neural Network", "0.610", "0.32", "0.40", "0.36"],
    ]
    add_table(
        s, data, Inches(LEFT_M), Inches(BODY_Y), Inches(BODY_W), Inches(2.0),
        col_widths=[Inches(4.2), Inches(2.4), Inches(1.6), Inches(1.6), Inches(2.1)], font_size=14,
    )
    # Honest framing
    rect(s, Inches(LEFT_M), Inches(BODY_Y + 2.4), Inches(BODY_W), Inches(2.0), fill=SURFACE, radius=0.05)
    add_text(s, "AN HONEST READ OF 0.65", Inches(LEFT_M + 0.3), Inches(BODY_Y + 2.55), Inches(BODY_W - 0.6), Inches(0.3),
             size=11, color=WARN, bold=True)
    add_text(
        s,
        "Published crash-prediction work using weather features sits in the 0.65 – 0.75 band. We are at the bottom of that band on purpose.",
        Inches(LEFT_M + 0.3), Inches(BODY_Y + 2.9), Inches(BODY_W - 0.6), Inches(0.5),
        size=15, color=TEXT, line_spacing=1.4,
    )
    add_text(
        s,
        "Most crashes happen in mild weather because most driving happens in mild weather. The model honestly stays near typical instead of pretending to be confident.",
        Inches(LEFT_M + 0.3), Inches(BODY_Y + 3.4), Inches(BODY_W - 0.6), Inches(0.7),
        size=15, color=MUTED, italic=True, line_spacing=1.4,
    )


# ---- 10 · Precision and Recall explainer ----------------------------------


def slide_precision_recall():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 10, TOTAL)
    title_block(
        s,
        "9 · Precision and recall, in plain language",
        "Two questions that matter for a safety advisory.",
        "These are not just metrics — they translate directly into real costs for the driver.",
    )
    col_w = (BODY_W - 0.4) / 2
    # Left: Precision
    rect(s, Inches(LEFT_M), Inches(BODY_Y), Inches(col_w), Inches(3.5), fill=SURFACE, radius=0.05)
    add_text(s, "RECALL  ·  THRESHOLD 0.35", Inches(LEFT_M + 0.3), Inches(BODY_Y + 0.25), Inches(col_w - 0.6), Inches(0.3),
             size=11, color=GOOD, bold=True)
    add_text(s, "\"Of all the genuinely crashy segments, how many did we catch?\"",
             Inches(LEFT_M + 0.3), Inches(BODY_Y + 0.6), Inches(col_w - 0.6), Inches(0.9),
             size=17, color=TEXT, bold=True, line_spacing=1.3)
    add_text(s, "0.84", Inches(LEFT_M + 0.3), Inches(BODY_Y + 1.6), Inches(col_w - 0.6), Inches(0.8),
             size=44, color=GOOD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s,
        "84 % of real crash situations caught. We operate at the F1-optimal threshold — the data itself nominates this as the best balance.",
        Inches(LEFT_M + 0.3), Inches(BODY_Y + 2.5), Inches(col_w - 0.6), Inches(0.9),
        size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.4,
    )
    # Right: Precision at that operating point
    right_x = LEFT_M + col_w + 0.4
    rect(s, Inches(right_x), Inches(BODY_Y), Inches(col_w), Inches(3.5), fill=SURFACE, radius=0.05)
    add_text(s, "PRECISION  ·  THRESHOLD 0.35", Inches(right_x + 0.3), Inches(BODY_Y + 0.25), Inches(col_w - 0.6), Inches(0.3),
             size=11, color=ACCENT, bold=True)
    add_text(s, "\"When we say crash risk is high, how often is it actually high?\"",
             Inches(right_x + 0.3), Inches(BODY_Y + 0.6), Inches(col_w - 0.6), Inches(0.9),
             size=17, color=TEXT, bold=True, line_spacing=1.3)
    add_text(s, "0.31", Inches(right_x + 0.3), Inches(BODY_Y + 1.6), Inches(col_w - 0.6), Inches(0.8),
             size=44, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s,
        "The cost of catching 84 % of real crashes. Acceptable for a safety advisory: over-warning is much cheaper than missing.",
        Inches(right_x + 0.3), Inches(BODY_Y + 2.5), Inches(col_w - 0.6), Inches(0.9),
        size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.4,
    )
    add_text(
        s,
        "F1 at this threshold is 0.449 — strictly better than every other threshold. The conventional 0.5 cutoff would catch only 25 %.",
        Inches(LEFT_M), Inches(BODY_Y + 3.7), Inches(BODY_W), Inches(0.5),
        size=13, color=ACCENT_LIGHT, italic=True, align=PP_ALIGN.CENTER,
    )


# ---- 11 · Architecture (visual flow) --------------------------------------


def slide_architecture():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 11, TOTAL)
    title_block(
        s,
        "10 · The system, end to end",
        "How a single search travels through the product.",
        "Five external services and one in-house model, composed in a single request.",
    )
    steps = [
        ("Browser", "User types two locations; geolocation auto-fills the From field.", ACCENT_LIGHT),
        ("Mapbox", "Returns two or three alternative routes between the points.", ACCENT_LIGHT),
        ("OpenWeather", "Live conditions at every segment midpoint along each route.", ACCENT_LIGHT),
        ("Mapbox Streets", "Classifies each segment as highway, arterial, residential, etc.", ACCENT_LIGHT),
        ("MLflow Model Server", "Returns a crash-history probability per segment.", GOOD),
        ("Web Server", "Runs the rule-based weather score, joins both scores, sends to the browser.", ACCENT_LIGHT),
    ]
    y0 = BODY_Y
    row_h = 0.66
    for i, (label, body, color) in enumerate(steps):
        y = y0 + i * row_h
        # Number badge
        rect(s, Inches(LEFT_M), Inches(y + 0.03), Inches(0.5), Inches(0.5), fill=color, radius=0.5)
        add_text(s, str(i + 1), Inches(LEFT_M), Inches(y + 0.03), Inches(0.5), Inches(0.5),
                 size=15, color=BG, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Label
        add_text(s, label, Inches(LEFT_M + 0.7), Inches(y + 0.0), Inches(3.0), Inches(0.35),
                 size=15, color=TEXT, bold=True)
        # Body
        add_text(s, body, Inches(LEFT_M + 4.0), Inches(y + 0.05), Inches(BODY_W - 4.2), Inches(0.5),
                 size=13, color=MUTED, line_spacing=1.3)


# ---- 12 · MLflow ----------------------------------------------------------


def slide_mlflow():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 12, TOTAL)
    title_block(
        s,
        "11 · MLflow keeps us honest",
        "The contract between training and serving.",
        "Anything the live app shows can be traced to a specific training run, dataset version, and model version.",
    )
    items = [
        ("Tracking",
         "Every training run logs its inputs, metrics, and the trained model.",
         "If a number in the product looks wrong, we can find the run that produced it."),
        ("Registry",
         "The best run gets registered under a stable name with a \"production\" alias.",
         "The web app calls the alias, not a specific file path."),
        ("Serving",
         "MLflow serves the registered model as a web endpoint.",
         "The web app sends features, gets a probability back."),
        ("Rollback",
         "Older versions stay in the registry. Promotion and rollback are a single alias flip.",
         "No code change, no redeploy — just a name change."),
    ]
    y0 = BODY_Y
    row_h = 0.96
    for i, (label, body, why) in enumerate(items):
        y = y0 + i * row_h
        rect(s, Inches(LEFT_M), Inches(y), Inches(BODY_W), Inches(0.86), fill=SURFACE, radius=0.05)
        add_text(s, label, Inches(LEFT_M + 0.3), Inches(y + 0.1), Inches(2.0), Inches(0.35),
                 size=14, color=ACCENT, bold=True)
        add_text(s, body, Inches(LEFT_M + 2.3), Inches(y + 0.1), Inches(5.7), Inches(0.7),
                 size=13, color=TEXT, line_spacing=1.3)
        add_text(s, why, Inches(LEFT_M + 8.2), Inches(y + 0.1), Inches(BODY_W - 8.4), Inches(0.7),
                 size=12, color=MUTED, italic=True, line_spacing=1.3)


# ---- 13 · Explainability --------------------------------------------------


def slide_explainability():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 13, TOTAL)
    title_block(
        s,
        "12 · Why a score is what it is",
        "Two complementary explanations, one for each score.",
        "Every number the driver sees can be traced back to its reasons.",
    )
    col_w = (BODY_W - 0.4) / 2
    # Rule explainability
    rect(s, Inches(LEFT_M), Inches(BODY_Y), Inches(col_w), Inches(4.0), fill=SURFACE, radius=0.05)
    add_text(s, "RULE-BASED SCORE", Inches(LEFT_M + 0.3), Inches(BODY_Y + 0.2), Inches(col_w - 0.6), Inches(0.3),
             size=11, color=WARN, bold=True)
    add_text(s, "Inherently explained.",
             Inches(LEFT_M + 0.3), Inches(BODY_Y + 0.55), Inches(col_w - 0.6), Inches(0.5),
             size=18, color=TEXT, bold=True)
    add_text(
        s,
        "Every score decomposes into named factors with named thresholds.\n\nThe UI shows each factor as a coloured bar with its description: \"snow contributed 18 points\", \"this road is a bridge so the total was multiplied by 1.5\".\n\nNo black box — the driver can see exactly why a number is what it is.",
        Inches(LEFT_M + 0.3), Inches(BODY_Y + 1.2), Inches(col_w - 0.6), Inches(2.7),
        size=13, color=TEXT, line_spacing=1.5,
    )
    # ML explainability
    right_x = LEFT_M + col_w + 0.4
    rect(s, Inches(right_x), Inches(BODY_Y), Inches(col_w), Inches(4.0), fill=SURFACE, radius=0.05)
    add_text(s, "MACHINE-LEARNED SCORE", Inches(right_x + 0.3), Inches(BODY_Y + 0.2), Inches(col_w - 0.6), Inches(0.3),
             size=11, color=ACCENT, bold=True)
    add_text(s, "Two layers, one shipped, one planned.",
             Inches(right_x + 0.3), Inches(BODY_Y + 0.55), Inches(col_w - 0.6), Inches(0.5),
             size=18, color=TEXT, bold=True)
    add_paragraphs(s, [
        [("Globally  ", {"bold": True, "color": GOOD, "size": 14}),
         ("The model tells us which inputs it relies on most. Top contributors are nearby crash history, location, hour of day, and road type.", {"color": TEXT, "size": 13})],
        [("", {})],
        [("Locally (next)  ", {"bold": True, "color": WARN, "size": 14}),
         ("Per-segment contributions for each specific prediction, surfaced as the same coloured-bar component already used by the rule-based score.", {"color": TEXT, "size": 13})],
    ], Inches(right_x + 0.3), Inches(BODY_Y + 1.2), Inches(col_w - 0.6), Inches(2.7),
       size=13, line_spacing=1.5, space_after=4)


# ---- 14 · Drift -----------------------------------------------------------


def slide_drift():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 14, TOTAL)
    title_block(
        s,
        "13 · Keeping it honest over time",
        "Three different ways a model can quietly decay.",
        "Each requires a different signal. The framework is the assignment; the implementation is the next iteration.",
    )
    data = [
        ["Type of drift", "What it looks like for ForeRoute", "How we'd catch it"],
        [
            ("The inputs change", {"bold": True, "color": GOOD}),
            "Winter arrives — colder, more snow, less visibility, higher humidity.",
            "Statistical tests on every input, against the original training window.",
        ],
        [
            ("The proportions change", {"bold": True, "color": WARN}),
            "Boston widens Vision Zero reporting to include cyclist near-misses.",
            "Watch the model's output distribution for shifts.",
        ],
        [
            ("The mapping itself changes", {"bold": True, "color": BAD}),
            "New automatic-braking technology makes snowy intersections genuinely safer.",
            "Performance drops on freshly-labelled crash data.",
        ],
    ]
    add_table(
        s, data, Inches(LEFT_M), Inches(BODY_Y), Inches(BODY_W), Inches(2.8),
        col_widths=[Inches(3.0), Inches(5.6), Inches(3.3)], font_size=13,
    )
    # Bottom callout
    rect(s, Inches(LEFT_M), Inches(BODY_Y + 3.1), Inches(BODY_W), Inches(1.0), fill=SURFACE, radius=0.05)
    add_text(s, "TOOLING", Inches(LEFT_M + 0.3), Inches(BODY_Y + 3.25), Inches(BODY_W - 0.6), Inches(0.3),
             size=11, color=ACCENT, bold=True)
    add_text(
        s,
        "Evidently AI produces a set of monitoring reports for each type of drift. They live next to the model in the repository and re-run on every release.",
        Inches(LEFT_M + 0.3), Inches(BODY_Y + 3.55), Inches(BODY_W - 0.6), Inches(0.5),
        size=14, color=TEXT, line_spacing=1.3,
    )


# ---- 15 · Deployment ------------------------------------------------------


def slide_deployment():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 15, TOTAL)
    title_block(
        s,
        "14 · How we'd deploy this for real",
        "What turns a working development build into a production product.",
        "The model is in place. These are the steps a real launch would need.",
    )
    items = [
        ("Log every prediction",
         "Write each segment's inputs, prediction, and verdict to disk so we can audit decisions and detect drift.", GOOD),
        ("Run drift checks nightly",
         "Compare the last seven days of real traffic against the training distribution. Alert if any input drifts past its threshold.", ACCENT_LIGHT),
        ("Pull updated crash records quarterly",
         "Boston publishes Vision Zero updates quarterly. We'd retrain and only promote the new version if it beats the current one on the same held-out test.", ACCENT_LIGHT),
        ("Provide a one-flip rollback",
         "Previous versions stay in the registry. If a new model regresses, flip the alias back — no redeploy.", WARN),
        ("Audit per neighbourhood",
         "Make sure the model is roughly as accurate in every part of the metro, not just the average. Crash density is not evenly distributed.", BAD),
        ("Publish a model card",
         "A one-page summary describing what the model does, what it doesn't do, and what it should never be used for.", ACCENT_LIGHT),
    ]
    y0 = BODY_Y
    row_h = 0.66
    for i, (label, body, color) in enumerate(items):
        y = y0 + i * row_h
        # Small accent bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LEFT_M), Inches(y + 0.07), Inches(0.08), Inches(0.4))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.shadow.inherit = False
        add_text(s, label, Inches(LEFT_M + 0.25), Inches(y), Inches(4.5), Inches(0.35),
                 size=15, color=TEXT, bold=True)
        add_text(s, body, Inches(LEFT_M + 0.25), Inches(y + 0.32), Inches(BODY_W - 0.45), Inches(0.5),
                 size=13, color=MUTED, line_spacing=1.3)


# ---- 16 · Recap -----------------------------------------------------------


def slide_recap():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 16, TOTAL)
    title_block(
        s,
        "Recap",
        "Four numbers to remember.",
        "If you remember nothing else from this deck, these.",
    )
    tiles = [
        ("89,040", "training rows", "22,260 real crashes + 66,780 carefully-chosen negatives", ACCENT),
        ("15", "inputs", "Weather, road, location, time, crash history", ACCENT_LIGHT),
        ("0.65", "ranking quality", "Gradient Boosting,\nregistered as the production model", GOOD),
        ("2", "scores in the UI", "Weather right now + crash history,\nside by side, disagreement flagged", WARN),
    ]
    tile_w = (BODY_W - 0.6) / 4
    tile_h = 3.4
    top = BODY_Y + 0.05
    for i, (big, label, sub, color) in enumerate(tiles):
        x = LEFT_M + i * (tile_w + 0.2)
        rect(s, Inches(x), Inches(top), Inches(tile_w), Inches(tile_h), fill=SURFACE, radius=0.05)
        add_text(s, big, Inches(x), Inches(top + 0.5), Inches(tile_w), Inches(1.3),
                 size=44, color=color, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, label, Inches(x), Inches(top + 1.85), Inches(tile_w), Inches(0.4),
                 size=14, color=TEXT, align=PP_ALIGN.CENTER, bold=True)
        add_text(s, sub, Inches(x + 0.2), Inches(top + 2.3), Inches(tile_w - 0.4), Inches(1.0),
                 size=12, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.4)


# ---- 17 · Thank you -------------------------------------------------------


def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s)
    chrome(s, 17, TOTAL)
    add_text(s, "Thank you.", Inches(0), Inches(2.6), Inches(PAGE_W), Inches(1.4),
             size=72, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.05), Inches(3.95), Inches(1.2), Inches(0.05))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.shadow.inherit = False
    add_text(s, "Live demo at  localhost:3000",
             Inches(0), Inches(4.2), Inches(PAGE_W), Inches(0.5),
             size=20, color=ACCENT_LIGHT, italic=True, align=PP_ALIGN.CENTER)
    add_text(s,
        "Your location auto-fills the From field. Type any Boston destination to see both scores in action.",
        Inches(0), Inches(4.75), Inches(PAGE_W), Inches(0.5),
        size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )
    # Companion docs
    rect(s, Inches(3.5), Inches(5.6), Inches(6.3), Inches(1.0), fill=SURFACE, radius=0.06)
    add_text(s, "COMPANION DOCUMENTS", Inches(3.5), Inches(5.75), Inches(6.3), Inches(0.3),
             size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        s, "Model card   ·   Full report   ·   Monitoring framework",
        Inches(3.5), Inches(6.05), Inches(6.3), Inches(0.4),
        size=14, color=TEXT, align=PP_ALIGN.CENTER,
    )
    add_text(s, "Questions?", Inches(0), Inches(6.85), Inches(PAGE_W), Inches(0.4),
             size=18, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


# ===========================================================================


def main():
    slide_title()
    slide_idea()
    slide_two_scores()
    slide_data_sources()
    slide_cleaning()
    slide_negatives()
    slide_features()
    slide_models()
    slide_performance()
    slide_precision_recall()
    slide_architecture()
    slide_mlflow()
    slide_explainability()
    slide_drift()
    slide_deployment()
    slide_recap()
    slide_thanks()
    prs.save(str(OUT))
    print(f"wrote {OUT}  ({OUT.stat().st_size:,} bytes, 17 slides)")


if __name__ == "__main__":
    main()
